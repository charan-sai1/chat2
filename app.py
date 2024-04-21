import google.generativeai as genai
import streamlit as st
import os


api = "AIzaSyB7R_31L1M5H3Qjwqd1LKy3QrHPM2zMbZM"
genai.configure(api_key=api)


model = genai.GenerativeModel(model_name="gemini-pro")


from PyPDF2 import PdfReader
from pptx import Presentation
import docx

def read_text_from_pdf(file):
    # Create a PDF reader object
    pdf_reader = PdfReader(file)
    
    # Initialize an empty string to store the concatenated content
    full_text = ""
    
    # Loop through each page of the PDF
    for page in pdf_reader.pages:
        # Extract text from the current page
        page_text = page.extract_text()
        
        # Append the page text to the full text string
        full_text += page_text
    
    return full_text

def read_text_from_docx(file):
    # Create a DOCX document object
    doc = docx.Document(file)
    
    # Initialize an empty string to store the concatenated content
    full_text = ""
    
    # Loop through each paragraph in the document
    for paragraph in doc.paragraphs:
        # Append the paragraph text to the full text string
        full_text += paragraph.text + "\n"
    
    return full_text

def read_text_from_ppt(file):
    # Create a PowerPoint presentation object
    ppt = Presentation(file)
    
    # Initialize an empty string to store the concatenated content
    full_text = ""
    
    # Loop through each slide of the presentation
    for slide in ppt.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text"):
                # Append the text to the full text string
                full_text += shape.text
    
    return full_text

def read_text(file):
    file_extension = file.name.split(".")[-1].lower()
    if file_extension == "pdf":
        return read_text_from_pdf(file)
    elif file_extension == "docx":
        return read_text_from_docx(file)
    elif file_extension == "pptx":
        return read_text_from_ppt(file)
    else:
        return "Unsupported file format"

def ask(question:str):
    try:
        if question!='':
            response = chat.send_message(content=question)
    
    except "invalid_grant" as e:
        token_path = f'\\token.json'
        if os.path.exists(token_path):
           os.remove(token_path)
        print("\n\n\n\!!!!!!!!!!!!!!!!!!!!!!!!!!!!!\n\n\n")
        return 
           
    except Exception as e:
        return
    return response.text

import streamlit as st
import os

if __name__ == "__main__":
    st.title("Chat with pdf")
    
    # File upload section
    st.header("Upload a File")
    uploaded_file = st.file_uploader("Choose a file", type=['PDF', 'doc'])
    if uploaded_file is not None:
        pdf_content = read_text(uploaded_file)
        if pdf_content:
            hist = [
                    {"parts": [{"text": pdf_content}], "role": "user"},
                    {"parts": [{"text": "okay. what do you want me to do ?"}], "role": "model"},
                    {"parts": [{"text": "now with the give data answer my questions in a detiled manner.strictly do not give the information which is beyond the given data. Help me to prepare for my exams explain all the things i should know related to the question asked."}], "role": "user"},
                    {"parts": [{"text": "shure. i wll not give any information beyond the give context.i will defenitely help you for your exams."}], "role": "model"},
            ]
            chat = model.start_chat(history=hist)
        else:
            st.write("Invalid file")    
    # Text input section
    st.header("Enter Text")
    text_input = st.text_input("Enter text here:")
    if st.button("Submit"):
        st.write(pdf_content)
